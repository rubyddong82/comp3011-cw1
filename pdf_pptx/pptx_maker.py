#!/usr/bin/env python3

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


SLIDE_WIDTH = 13.333  # inches, default widescreen
SLIDE_HEIGHT = 7.5


def set_slide_background(slide, color=(245, 247, 250)):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*color)


def add_title(slide, text):
    box = slide.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(12.1), Inches(0.8))
    tf = box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(20, 40, 80)
    p.alignment = PP_ALIGN.CENTER
    return box


def add_subtitle(slide, text):
    box = slide.shapes.add_textbox(Inches(0.8), Inches(1.05), Inches(11.7), Inches(0.45))
    tf = box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(14)
    run.font.color.rgb = RGBColor(90, 90, 90)
    p.alignment = PP_ALIGN.CENTER
    return box


def add_content_box(slide, lines, left=1.0, top=1.8, width=11.3, height=4.5):
    shape = slide.shapes.add_shape(
        1,  # MSO_AUTO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )

    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)

    line = shape.line
    line.color.rgb = RGBColor(200, 205, 210)

    tf = shape.text_frame
    tf.clear()

    for i, line_text in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line_text
        p.level = 0
        p.alignment = PP_ALIGN.LEFT

        for run in p.runs:
            run.font.size = Pt(20)
            run.font.color.rgb = RGBColor(25, 25, 25)

    return shape


def add_footer(slide, text):
    box = slide.shapes.add_textbox(Inches(0.5), Inches(6.85), Inches(12.3), Inches(0.3))
    tf = box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(10)
    run.font.color.rgb = RGBColor(120, 120, 120)
    p.alignment = PP_ALIGN.RIGHT
    return box


def add_link_list(slide, title, items, left=1.0, top=2.0, width=11.3, height=3.8):
    shape = slide.shapes.add_shape(
        1,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )

    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)

    line = shape.line
    line.color.rgb = RGBColor(200, 205, 210)

    tf = shape.text_frame
    tf.clear()

    p0 = tf.paragraphs[0]
    p0.text = title
    for run in p0.runs:
        run.font.size = Pt(22)
        run.font.bold = True
        run.font.color.rgb = RGBColor(20, 40, 80)

    for label, url in items:
        p = tf.add_paragraph()
        p.text = f"{label}: {url}"
        for run in p.runs:
            run.font.size = Pt(18)
            run.font.color.rgb = RGBColor(30, 30, 30)

    return shape


def make_slide(prs, title, subtitle, bullet_lines, footer):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_background(slide)
    add_title(slide, title)
    add_subtitle(slide, subtitle)
    add_content_box(slide, [f"• {x}" for x in bullet_lines])
    add_footer(slide, footer)
    return slide


def build_demo_pptx(output_path="demo_submission_slides.pptx"):
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_WIDTH)
    prs.slide_height = Inches(SLIDE_HEIGHT)

    # remove default first slide if present by recreating fresh usage pattern
    # python-pptx starts empty in Presentation(), so nothing extra needed

    # 1. Intro
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    add_title(slide, "Project Submission Presentation")
    add_subtitle(slide, "Demo deck generated with python-pptx")
    add_content_box(slide, [
        "• Project title, student name, module code",
        "• Brief overview of the submission package",
        "• This demo has 6 slides total",
        "• Replace this content with your real presentation text"
    ])
    add_footer(slide, "Slide 1 of 6")

    # 2. Version control
    make_slide(
        prs,
        "Version Control Practices and Commit History",
        "How the development process was managed in Git",
        [
            "Repository structure and branch usage",
            "Commit frequency across development stages",
            "Examples of meaningful commit messages",
            "Evidence of iterative implementation, debugging, and refinement"
        ],
        "Slide 2 of 6"
    )

    # 3. API docs
    make_slide(
        prs,
        "API Documentation Overview",
        "What is documented and how users can navigate it",
        [
            "Purpose of the API and core endpoint groups",
            "Where the documentation is hosted",
            "Request / response format and usage examples",
            "Setup, testing, and developer guidance"
        ],
        "Slide 3 of 6"
    )

    # 4. Technical report
    make_slide(
        prs,
        "Technical Report Highlights",
        "Main points from the written report",
        [
            "Architecture and implementation summary",
            "Important design decisions and trade-offs",
            "Testing, validation, and limitations",
            "GenAI declaration and appendix coverage"
        ],
        "Slide 4 of 6"
    )

    # 5. Deliverables
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    add_title(slide, "All Deliverables")
    add_subtitle(slide, "Keep all final submission assets in one place")
    add_link_list(
        slide,
        "Submission Links",
        [
            ("Code repository", "https://github.com/yourname/yourrepo"),
            ("API documentation", "https://github.com/yourname/yourrepo/docs"),
            ("Technical report", "https://example.com/report.pdf"),
            ("Presentation slides", "https://example.com/slides.pptx"),
        ],
        top=1.9,
        height=4.2
    )
    add_footer(slide, "Slide 5 of 6")

    # 6. Ending
    make_slide(
        prs,
        "Thank You",
        "Closing slide",
        [
            "Final recap of repository, documentation, report, and slides",
            "Confirm that all required deliverables are included",
            "Add contact details if needed",
            "Q&A / submission complete"
        ],
        "Slide 6 of 6"
    )

    prs.save(output_path)
    print(output_path)


if __name__ == "__main__":
    build_demo_pptx()